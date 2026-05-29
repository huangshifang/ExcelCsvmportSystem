"""
Generate promotional PPT for Excel Import System.
Requires: python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Color scheme
PRIMARY = RGBColor(0x1A, 0x56, 0xDB)      # Blue
SECONDARY = RGBColor(0x10, 0x98, 0xAD)    # Teal
ACCENT = RGBColor(0x25, 0xAE, 0x5F)       # Green
DARK = RGBColor(0x1E, 0x29, 0x3B)         # Dark blue-gray
LIGHT = RGBColor(0xF3, 0xF4, 0xF6)        # Light gray
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
GRAY = RGBColor(0x6B, 0x72, 0x80)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xEF, 0x44, 0x44)
DARK_BG = RGBColor(0x0F, 0x17, 0x2A)      # Very dark blue

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLD_W = Inches(13.333)
SLD_H = Inches(7.5)


def add_bg(slide, color):
    """Set solid background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color, opacity=None):
    """Add a colored rectangle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if opacity is not None:
        shape.fill.fore_color.brightness = opacity
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=BLACK, bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    """Add a text box with single text run."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multi_text(slide, left, top, width, height, lines, font_size=16,
                   color=BLACK, line_spacing=1.5, font_name='Microsoft YaHei', alignment=PP_ALIGN.LEFT):
    """Add text box with multiple paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if isinstance(line_data, str):
            p.text = line_data
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
            p.font.name = font_name
        else:
            p.text = line_data.get('text', '')
            p.font.size = Pt(line_data.get('size', font_size))
            p.font.color.rgb = line_data.get('color', color)
            p.font.bold = line_data.get('bold', False)
            p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(font_size * (line_spacing - 1))
    return txBox


def add_card(slide, left, top, width, height, title, lines, title_color=PRIMARY, bg_color=WHITE):
    """Add a card-style box with title and content."""
    shape = add_rect(slide, left, top, width, height, bg_color)
    shape.shadow.inherit = False
    # Title
    add_text_box(slide, left + Inches(0.3), top + Inches(0.15), width - Inches(0.6), Inches(0.5),
                 title, font_size=18, color=title_color, bold=True)
    # Content
    add_multi_text(slide, left + Inches(0.3), top + Inches(0.65), width - Inches(0.6),
                   height - Inches(0.8), lines, font_size=13, color=DARK, line_spacing=1.4)


def add_title_bar(slide, title, subtitle=None):
    """Add a top title bar."""
    add_rect(slide, Inches(0), Inches(0), SLD_W, Inches(1.1), DARK)
    add_text_box(slide, Inches(0.8), Inches(0.15), Inches(11), Inches(0.7),
                 title, font_size=30, color=WHITE, bold=True)
    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.4),
                     subtitle, font_size=14, color=RGBColor(0x9C, 0xA3, 0xAF))


def add_footer(slide, page_num):
    """Add a footer with page number."""
    add_rect(slide, Inches(0), SLD_H - Inches(0.4), SLD_W, Inches(0.4), DARK)
    add_text_box(slide, Inches(0.5), SLD_H - Inches(0.35), Inches(6), Inches(0.3),
                 'Excel Import System | 企业级数据导入平台', font_size=9, color=GRAY)
    add_text_box(slide, Inches(11.5), SLD_H - Inches(0.35), Inches(1.5), Inches(0.3),
                 f'{page_num}', font_size=9, color=GRAY, alignment=PP_ALIGN.RIGHT)


# ============================================================
# SLIDE 1: TITLE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK_BG)

# Decorative shapes
add_rect(slide, Inches(0), Inches(0), SLD_W, Inches(0.08), PRIMARY)
add_rect(slide, Inches(0), SLD_H - Inches(0.08), SLD_W, Inches(0.08), SECONDARY)

# Large decorative circle
circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.5), Inches(1.5), Inches(6), Inches(6))
circle.fill.solid()
circle.fill.fore_color.rgb = PRIMARY
circle.fill.fore_color.brightness = 0.85
circle.line.fill.background()

circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.5), Inches(3.0), Inches(3.5), Inches(3.5))
circle2.fill.solid()
circle2.fill.fore_color.rgb = SECONDARY
circle2.fill.fore_color.brightness = 0.7
circle2.line.fill.background()

# Title text
add_text_box(slide, Inches(1.0), Inches(1.8), Inches(8), Inches(0.8),
             'Excel Import System', font_size=48, color=WHITE, bold=True)
add_text_box(slide, Inches(1.0), Inches(2.7), Inches(8), Inches(0.6),
             '企业级 Excel / CSV 数据导入平台', font_size=28, color=SECONDARY, bold=False)

# Divider line
add_rect(slide, Inches(1.0), Inches(3.5), Inches(2.5), Inches(0.06), ACCENT)

add_multi_text(slide, Inches(1.0), Inches(3.9), Inches(7), Inches(1.5), [
    {'text': '.NET 10 + React 19 + SQL Server | Clean Architecture', 'size': 18, 'color': RGBColor(0x9C, 0xA3, 0xAF)},
    {'text': 'JWT 混合认证 (本地 + LDAP/AD) | RBAC 角色权限 | 多服务器支持', 'size': 16, 'color': GRAY},
], line_spacing=1.8)

add_text_box(slide, Inches(1.0), Inches(5.5), Inches(5), Inches(0.4),
             '2026年5月  |  Version 1.4.0', font_size=14, color=GRAY)


# ============================================================
# SLIDE 2: PROJECT OVERVIEW
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, '项目背景与定位', 'PROJECT OVERVIEW')
add_footer(slide, '02 / 14')

# Problem statement
add_card(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(2.5),
         '痛点分析', [
             '• 业务人员依赖 DBA 手动执行 SQL 导入数据，效率低',
             '• Excel/CSV 数据格式多样，列映射易出错',
             '• 多 SQL Server 实例分散管理，缺乏统一导入入口',
             '• 缺乏细粒度的表级权限控制和操作审计',
         ], title_color=RED)

add_card(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(2.5),
         '解决方案', [
             '• Web 化自助导入：业务人员浏览器操作，无需数据库客户端',
             '• 智能列映射：自动匹配 Excel 列与数据库列，支持手动调整',
             '• 多服务器统一管理：注册远程实例，跨服务器导入',
             '• 企业级安全：RBAC + 表级权限 + 完整审计日志',
         ], title_color=ACCENT)

# Key numbers
for i, (num, label) in enumerate([
    ('3层', 'Clean Architecture\n架构分层'),
    ('10+', '安全防护\n机制'),
    ('3种', '用户角色\nAdmin/Op/Viewer'),
    ('200MB', '大文件\n上传支持'),
]):
    x = Inches(0.5 + i * 3.2)
    add_rect(slide, x, Inches(4.6), Inches(2.9), Inches(1.1), LIGHT)
    add_text_box(slide, x + Inches(0.2), Inches(4.7), Inches(2.5), Inches(0.5),
                 num, font_size=28, color=PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), Inches(5.2), Inches(2.5), Inches(0.4),
                 label, font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)

# Bottom summary
add_text_box(slide, Inches(0.5), Inches(6.2), Inches(12), Inches(0.5),
             '定位：企业级、安全、可扩展的数据导入中间件，打通 Excel → SQL Server 的"最后一公里"',
             font_size=14, color=DARK, bold=False, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 3: TECH STACK
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, '核心技术栈', 'TECHNOLOGY STACK')
add_footer(slide, '03 / 14')

stacks = [
    ('前端', PRIMARY, [
        'React 19 + TypeScript',
        'Vite 8 构建工具',
        'Ant Design 6 组件库',
        'React Router 7 路由',
        'i18next 国际化 (中/英)',
        'Axios HTTP 客户端',
    ]),
    ('后端', SECONDARY, [
        '.NET 10 (ASP.NET Core)',
        'Entity Framework Core 10',
        'JWT Bearer 认证',
        'BCrypt 密码哈希',
        'EPPlus 7 (Excel读写)',
        'CsvHelper 33 (CSV解析)',
    ]),
    ('基础设施', ACCENT, [
        'SQL Server 数据库',
        'Docker 容器化部署',
        'Nginx HTTPS 反向代理',
        'Serilog 日志 (Console+File)',
        'Clean Architecture 架构',
        '自建内部 CA 证书体系',
    ]),
    ('安全体系', ORANGE, [
        'SVG 验证码 (防机器人)',
        '速率限制 (10次/分钟)',
        '账号锁定 (5次→15分钟)',
        'LDAP Filter 注入防护',
        'SQL 参数化查询',
        '登录审计全量记录',
    ]),
]

for i, (title, color, items) in enumerate(stacks):
    x = Inches(0.4 + i * 3.2)
    add_rect(slide, x, Inches(1.5), Inches(3.0), Inches(0.55), color)
    add_text_box(slide, x, Inches(1.52), Inches(3.0), Inches(0.5),
                 title, font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_rect(slide, x, Inches(2.05), Inches(3.0), Inches(3.3), LIGHT)
    add_multi_text(slide, x + Inches(0.2), Inches(2.2), Inches(2.6), Inches(3.0),
                   [f'✓ {item}' for item in items], font_size=13, color=DARK, line_spacing=1.8)


# ============================================================
# SLIDE 4: ARCHITECTURE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, '系统架构 — Clean Architecture 三层分离', 'ARCHITECTURE OVERVIEW')
add_footer(slide, '04 / 14')

# Architecture diagram
layers = [
    ('ExcelImportSystem.API', '表现层 (ASP.NET Core Web API)', PRIMARY,
     ['Controllers 处理 HTTP 请求/响应', 'JWT 认证 & CORS 配置', '种子数据 & 数据库迁移', '请求 → DTO → Service → 响应']),
    ('ExcelImportSystem.Infrastructure', '基础设施层 (EF Core + Services)', SECONDARY,
     ['AppDbContext 数据上下文', 'Auth / Import / Table / LDAP Service', 'DI 注册扩展 (AddInfrastructure)', 'IServiceScopeFactory 作用域管理']),
    ('ExcelImportSystem.Core', '领域层 (零外部依赖)', ACCENT,
     ['Entity 实体定义 (User/Role/ImportLog/...)', 'DTO 数据传输对象', 'Service 接口 (IAuthService/...)', 'Config Models (LdapSettings)']),
]

for i, (name, desc, color, items) in enumerate(layers):
    y = Inches(1.5 + i * 1.8)
    # Label box
    add_rect(slide, Inches(0.5), y, Inches(3.5), Inches(1.5), color)
    add_text_box(slide, Inches(0.7), y + Inches(0.15), Inches(3.1), Inches(0.45),
                 name, font_size=16, color=WHITE, bold=True)
    add_text_box(slide, Inches(0.7), y + Inches(0.55), Inches(3.1), Inches(0.35),
                 desc, font_size=11, color=RGBColor(0xE0, 0xE7, 0xFF))
    # Items
    add_rect(slide, Inches(4.3), y, Inches(8.5), Inches(1.5), LIGHT)
    add_multi_text(slide, Inches(4.6), y + Inches(0.15), Inches(7.9), Inches(1.2),
                   items, font_size=13, color=DARK, line_spacing=1.6)

# Arrows between layers
for i in range(2):
    y = Inches(3.0 + i * 1.8)
    arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(2.0), y, Inches(0.4), Inches(0.35))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = GRAY
    arrow.line.fill.background()

# Frontend note
add_rect(slide, Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.45), DARK)
add_text_box(slide, Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.45),
             '前端：React 19 SPA → Axios + JWT Interceptor → REST API → JSON 响应 (ApiResponse<T>)  |  部署：Docker Compose (API + Nginx + SQL Server)',
             font_size=11, color=WHITE, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 5: CORE FEATURES OVERVIEW
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, '核心功能全景', 'FEATURE OVERVIEW')
add_footer(slide, '05 / 14')

features = [
    ('📊\n数据导入', PRIMARY, 'Excel/CSV 智能导入\n两步流程 + 异步轮询\n自动列映射\n批次处理 (100~10000行)'),
    ('🔐\n认证安全', SECONDARY, 'JWT + BCrypt + LDAP\nSVG验证码防机器人\n速率限制 & 账号锁定\n登录审计全量记录'),
    ('👥\n权限管理', ACCENT, '三级角色 (A/O/V)\n8+1 细粒度权限\n表级数据库权限\nRBAC 前后端双控'),
    ('🖥️\n多服务器', ORANGE, '多 SQL Server 注册\n跨服务器数据导入\n连接测试 & 缓存\n优雅降级容错'),
    ('📋\n审计日志', RGBColor(0x8B, 0x5C, 0xF6), '导入操作日志\n登录审计日志\nIP/UA 全量记录\n分页搜索过滤'),
    ('🌐\n企业集成', RGBColor(0xEC, 0x48, 0x99), 'LDAP/AD 域认证\n运行时配置即时生效\nHTTPS + 内部 CA\nDocker 离线部署'),
]

for i, (name, color, desc) in enumerate(features):
    col = i % 3
    row = i // 3
    x = Inches(0.5 + col * 4.2)
    y = Inches(1.4 + row * 2.9)
    add_rect(slide, x, y, Inches(3.9), Inches(2.6), LIGHT)
    add_rect(slide, x, y, Inches(3.9), Inches(0.08), color)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.25), Inches(3.3), Inches(0.9),
                 name, font_size=17, color=color, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(1.3), Inches(3.3), Inches(1.1),
                 desc, font_size=12, color=DARK)


# ============================================================
# SLIDE 6: AUTH & SECURITY
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_rect(slide, Inches(0), Inches(0), SLD_W, Inches(1.1), RGBColor(0x0A, 0x0F, 0x1A))
add_text_box(slide, Inches(0.8), Inches(0.15), Inches(11), Inches(0.7),
             '认证与安全体系', font_size=30, color=WHITE, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.4),
             'AUTHENTICATION & SECURITY — 10+ 层安全防护', font_size=14, color=RGBColor(0x9C, 0xA3, 0xAF))
add_footer(slide, '06 / 14')

# Security flow
flow_items = [
    ('① CAPTCHA\n验证码', 'SVG 4位验证码\n每次登录必填\n5分钟过期\n自动刷新'),
    ('② Rate Limit\n速率限制', '固定窗口算法\n10次/分钟/IP\n超限返回 503\n框架内置保护'),
    ('③ Account\nLockout', '5次失败→锁定\n15分钟自动解锁\n防暴力猜解\n防用户名枚举'),
    ('④ JWT\nToken', 'BCrypt 密码哈希\nClaims 权限嵌入\n24小时过期\nBearer 认证'),
    ('⑤ LDAP\n认证', 'LDAP Filter 转义\nSSL/TLS 加密\n防注入攻击\n用完即释放'),
    ('⑥ Audit\n审计', '登录全量记录\nIP + UserAgent\n成功/失败追踪\n仅 Admin 可查'),
]

for i, (title, desc) in enumerate(flow_items):
    x = Inches(0.4 + i * 2.15)
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.7), Inches(1.4), Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = PRIMARY if i % 2 == 0 else SECONDARY
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(i + 1)
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    add_text_box(slide, x, Inches(2.2), Inches(2.0), Inches(0.6),
                 title, font_size=15, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(2.9), Inches(2.0), Inches(2.0),
                 desc, font_size=11, color=RGBColor(0x9C, 0xA3, 0xAF), alignment=PP_ALIGN.CENTER)

# Bottom highlight
add_rect(slide, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.8),
         RGBColor(0x15, 0x1F, 0x35))
add_text_box(slide, Inches(0.8), Inches(5.35), Inches(11.7), Inches(0.4),
             '🛡️  安全特性亮点', font_size=18, color=ACCENT, bold=True)
add_multi_text(slide, Inches(0.8), Inches(5.85), Inches(11.7), Inches(1.0), [
    {'text': '• CAPTCHA + 速率限制 + 账号锁定 三道防线防暴力破解    • 登录审计日志全量记录，支持安全事件追溯', 'size': 13, 'color': RGBColor(0xD1, 0xD5, 0xDB)},
    {'text': '• SVG 纯代码生成验证码，零原生依赖，Docker 跨平台兼容    • SQL 参数化查询 + LDAP Filter 转义，防止注入攻击', 'size': 13, 'color': RGBColor(0xD1, 0xD5, 0xDB)},
    {'text': '• 前端 autoComplete=off 防浏览器密码管理器覆盖    • 内部 CA 证书 + HTTPS 加密传输，防止中间人攻击', 'size': 13, 'color': RGBColor(0xD1, 0xD5, 0xDB)},
], line_spacing=1.4)


# ============================================================
# SLIDE 7: IMPORT FLOW
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, 'Excel / CSV 数据导入流程', 'IMPORT WORKFLOW — 两步流程 + 异步轮询')
add_footer(slide, '07 / 14')

# Step 1
add_rect(slide, Inches(0.5), Inches(1.5), Inches(3.8), Inches(0.6), PRIMARY)
add_text_box(slide, Inches(0.7), Inches(1.55), Inches(3.4), Inches(0.5),
             '① 上传与配置', font_size=20, color=WHITE, bold=True)
add_rect(slide, Inches(0.5), Inches(2.1), Inches(3.8), Inches(2.2), LIGHT)
add_multi_text(slide, Inches(0.7), Inches(2.2), Inches(3.4), Inches(2.0), [
    '• 选择目标服务器 (本地/远程)',
    '• 选择数据库和数据表',
    '• 上传 Excel/CSV 文件',
    '• 配置文件选项:',
    '  - 包含表头 (默认开启)',
    '  - 事务模式 (原子性保证)',
    '  - 批次大小 (100~10000)',
], font_size=13, color=DARK, line_spacing=1.5)

# Arrow
arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.5), Inches(2.2), Inches(0.5), Inches(0.4))
arrow.fill.solid()
arrow.fill.fore_color.rgb = GRAY
arrow.line.fill.background()

# Step 2
add_rect(slide, Inches(5.2), Inches(1.5), Inches(3.8), Inches(0.6), SECONDARY)
add_text_box(slide, Inches(5.4), Inches(1.55), Inches(3.4), Inches(0.5),
             '② 列映射与预览', font_size=20, color=WHITE, bold=True)
add_rect(slide, Inches(5.2), Inches(2.1), Inches(3.8), Inches(2.2), LIGHT)
add_multi_text(slide, Inches(5.4), Inches(2.2), Inches(3.4), Inches(2.0), [
    '• 自动读取文件列名和示例数据',
    '• 智能列映射 (不区分大小写)',
    '• 手动调整未映射的列',
    '• 跳过不需要导入的列',
    '• 自增列 / 主键自动跳过',
    '• 预览最多 5 行示例数据',
], font_size=13, color=DARK, line_spacing=1.5)

# Arrow
arrow2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(9.2), Inches(2.2), Inches(0.5), Inches(0.4))
arrow2.fill.solid()
arrow2.fill.fore_color.rgb = GRAY
arrow2.line.fill.background()

# Step 3
add_rect(slide, Inches(9.9), Inches(1.5), Inches(3.0), Inches(0.6), ACCENT)
add_text_box(slide, Inches(10.1), Inches(1.55), Inches(2.6), Inches(0.5),
             '③ Fire-and-Forget 异步执行', font_size=15, color=WHITE, bold=True)
add_rect(slide, Inches(9.9), Inches(2.1), Inches(3.0), Inches(2.2), LIGHT)
add_multi_text(slide, Inches(10.1), Inches(2.2), Inches(2.6), Inches(2.0), [
    '• 立即返回 taskId',
    '• 后台异步批量导入',
    '• 参数化 SQL 防注入',
    '• 批量插入 (默认1000行)',
    '• 可选事务保证原子性',
    '• 错误行收集与记录',
], font_size=12, color=DARK, line_spacing=1.5)

# Bottom section: progress polling
add_rect(slide, Inches(0.5), Inches(4.7), Inches(12.3), Inches(2.4), DARK)
add_text_box(slide, Inches(0.8), Inches(4.85), Inches(11.7), Inches(0.4),
             '④ 实时进度轮询 (每 500ms)', font_size=18, color=WHITE, bold=True)

# Status flow
statuses = ['pending\n等待中', 'reading\n读取文件', 'importing\n导入中', 'completed\n完成 ✓', 'failed\n失败 ✗']
for i, s in enumerate(statuses):
    x = Inches(1.0 + i * 2.4)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, Inches(5.4), Inches(0.8), Inches(0.8))
    circle.fill.solid()
    colors = [GRAY, ORANGE, PRIMARY, ACCENT, RED]
    circle.fill.fore_color.rgb = colors[i]
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = s
    p.font.size = Pt(10)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    if i < 4:
        arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + Inches(0.9), Inches(5.65), Inches(0.5), Inches(0.25))
        arr.fill.solid()
        arr.fill.fore_color.rgb = GRAY
        arr.line.fill.background()

add_multi_text(slide, Inches(0.8), Inches(6.35), Inches(11.7), Inches(0.6), [
    {'text': '技术要点: Fire-and-Forget 模式 | ConcurrentDictionary 内存进度追踪 | IServiceScopeFactory 作用域管理 | 最大 200MB 文件 | 动态目标表 INSERT',
     'size': 11, 'color': RGBColor(0x9C, 0xA3, 0xAF)}
], line_spacing=1.3)


# ============================================================
# SLIDE 8: MULTI-SERVER SUPPORT
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, '多 SQL Server 实例支持', 'MULTI-SERVER SUPPORT — 跨服务器数据导入')
add_footer(slide, '08 / 14')

# Architecture diagram
add_rect(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.0), LIGHT)
add_text_box(slide, Inches(0.8), Inches(1.6), Inches(5.2), Inches(0.4),
             '架构设计', font_size=18, color=PRIMARY, bold=True)
add_multi_text(slide, Inches(0.8), Inches(2.1), Inches(5.2), Inches(4.2), [
    {'text': 'IConnectionFactory (Singleton)', 'size': 15, 'bold': True, 'color': PRIMARY},
    '',
    '• serverId=null → 本地 SQL Server',
    '• serverId=N   → 远程 SQL Server 实例',
    '',
    '核心能力:',
    '• ConcurrentDictionary 缓存连接字符串',
    '• 懒加载数据库→服务器映射',
    '• ChangeDatabase() 上下文切换',
    '• 远程不可达时优雅降级 (仅 LogWarning)',
    '• INVALIDATE CACHE 缓存失效机制',
    '',
    '跨服务器操作:',
    '• 数据库发现 (sys.databases)',
    '• 表/列元数据查询 (INFORMATION_SCHEMA)',
    '• 数据导入 (SqlBulkCopy / Raw INSERT)',
    '• 表级权限验证 (含 ServerId)',
    '• ImportLog 记录 ServerId + ServerName',
], font_size=12, color=DARK, line_spacing=1.15)

# Right side: key patterns
add_card(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(2.2),
         '前端复合键模式', [
             '• 多服务器可能有同名数据库',
             '• 下拉选项使用复合键: ${serverId}::${dbName}',
             '• 解析: serverId=0 → undefined (本地)',
             '• 显示: [ServerName] DatabaseName',
             '• 表级权限树按服务器分组展示',
         ], title_color=SECONDARY)

add_card(slide, Inches(6.8), Inches(4.0), Inches(5.8), Inches(2.5),
         'ServerController API', [
             'GET    /api/server          — 列出所有服务器',
             'GET    /api/server/{id}     — 获取服务器详情',
             'POST   /api/server          — 注册新服务器 (ServerManage)',
             'PUT    /api/server/{id}     — 更新服务器配置',
             'DELETE /api/server/{id}     — 删除 (检查关联权限)',
             'POST   /api/server/test     — 测试连接字符串',
             '',
             'ServerId 在所有层级传播: DTO → Service → Controller → FormData',
         ], title_color=ACCENT)


# ============================================================
# SLIDE 9: RBAC & PERMISSIONS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, '权限控制体系 — RBAC + 表级权限', 'ROLE-BASED ACCESS CONTROL')
add_footer(slide, '09 / 14')

# Role table
roles = [
    ('Admin\n管理员', PRIMARY, ['Import.Execute', 'Import.View', 'User.Manage', 'Role.Manage',
                                 'Log.View', 'Audit.View', 'Database.Manage', 'System.Manage',
                                 'Server.View', 'Server.Manage']),
    ('Operator\n操作员', SECONDARY, ['Import.Execute', 'Import.View', 'Log.View']),
    ('Viewer\n查看者', ACCENT, ['Import.View']),
]

for i, (role, color, perms) in enumerate(roles):
    x = Inches(0.5 + i * 4.2)
    add_rect(slide, x, Inches(1.5), Inches(3.9), Inches(0.6), color)
    add_text_box(slide, x, Inches(1.52), Inches(3.9), Inches(0.55),
                 role, font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_rect(slide, x, Inches(2.1), Inches(3.9), Inches(1.7), LIGHT)
    perm_text = '\n'.join([f'  ✓ {p}' for p in perms])
    add_text_box(slide, x + Inches(0.2), Inches(2.15), Inches(3.5), Inches(1.6),
                 perm_text, font_size=11, color=DARK)

# Table-level permission
add_text_box(slide, Inches(0.5), Inches(4.1), Inches(12), Inches(0.5),
             '表级数据库权限 (Table-Level Access Control)', font_size=20, color=DARK, bold=True)

add_rect(slide, Inches(0.5), Inches(4.6), Inches(6.2), Inches(2.3), LIGHT)
add_multi_text(slide, Inches(0.8), Inches(4.7), Inches(5.6), Inches(2.1), [
    {'text': '权限粒度', 'size': 15, 'bold': True, 'color': PRIMARY},
    '',
    '• 整库授权 (Wildcard): 用户可导入该库所有表',
    '• 逐表授权 (Specific): 仅允许导入指定表',
    '• 按服务器独立管理 (含 ServerId 过滤)',
    '• 向后兼容: v1.0 整库权限自动转为 Wildcard',
    '',
    '• Admin 默认拥有全部数据库和表的访问权限',
    '• 非 Admin 用户仅能看到被授权的表',
], font_size=12, color=DARK, line_spacing=1.15)

add_rect(slide, Inches(7.0), Inches(4.6), Inches(5.8), Inches(2.3), LIGHT)
add_multi_text(slide, Inches(7.3), Inches(4.7), Inches(5.2), Inches(2.1), [
    {'text': '双重验证', 'size': 15, 'bold': True, 'color': SECONDARY},
    '',
    '前端 UI 控制 (hasPermission):',
    '  • 菜单项显示/隐藏',
    '  • 按钮启用/禁用',
    '  • 下拉选项过滤',
    '',
    '后端 API 验证 (Authorization Policy):',
    '  • Controller [Authorize] + Policy',
    '  • ImportService.ValidateAccess()',
    '  → Preview 和 Execute 均校验',
    '  → 即使绕过前端也无法越权',
], font_size=12, color=DARK, line_spacing=1.15)


# ============================================================
# SLIDE 10: LDAP/AD INTEGRATION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, 'LDAP / Active Directory 域集成', 'ENTERPRISE AD INTEGRATION')
add_footer(slide, '10 / 14')

add_rect(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(2.5), LIGHT)
add_text_box(slide, Inches(0.8), Inches(1.6), Inches(5.2), Inches(0.4),
             '认证流程', font_size=18, color=PRIMARY, bold=True)
add_multi_text(slide, Inches(0.8), Inches(2.1), Inches(5.2), Inches(1.8), [
    '1. 用户输入域账号 + 域密码',
    '2. 系统尝试本地 BCrypt 认证 → 失败',
    '3. 使用服务账号绑定 LDAP 服务器',
    '4. SearchRequest 查找用户 DN',
    '5. 使用用户 DN + 密码进行 Bind 验证',
    '6. 成功 → 自动创建本地用户 (Viewer 角色) → 签发 JWT',
    '7. 失败 → 记录失败次数，达阈值锁定账号',
], font_size=13, color=DARK, line_spacing=1.5)

add_rect(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(2.5), LIGHT)
add_text_box(slide, Inches(7.1), Inches(1.6), Inches(5.2), Inches(0.4),
             '技术实现', font_size=18, color=SECONDARY, bold=True)
add_multi_text(slide, Inches(7.1), Inches(2.1), Inches(5.2), Inches(1.8), [
    {'text': 'System.DirectoryServices.Protocols (S.DS.P)', 'size': 14, 'bold': True, 'color': SECONDARY},
    '',
    '• .NET 内置库，无需额外 NuGet 包',
    '• 跨平台: Windows / Linux / macOS',
    '• 支持 SSL/TLS 加密 (LDAPS :636)',
    '• LDAP Filter 转义防注入攻击',
    '• 连接用完即释放 (using 模式)',
    '• 服务账号密码支持环境变量注入',
], font_size=13, color=DARK, line_spacing=1.5)

# Configuration
add_text_box(slide, Inches(0.5), Inches(4.4), Inches(12), Inches(0.5),
             '运行时配置 — 即时生效，无需重启', font_size=18, color=DARK, bold=True)

add_rect(slide, Inches(0.5), Inches(4.9), Inches(12.3), Inches(2.2), DARK)
add_multi_text(slide, Inches(0.8), Inches(5.0), Inches(11.7), Inches(2.0), [
    {'text': 'appsettings.json (初始默认)  →  SystemSettings 表 (持久化)  →  LdapSettingsProvider (内存缓存)  →  LdapService',
     'size': 14, 'bold': True, 'color': ACCENT},
    '',
    {'text': '• 管理员可通过「系统设置」Web 页面修改 LDAP 配置，保存后即时生效，无需重启服务或修改配置文件', 'size': 13, 'color': RGBColor(0xD1, 0xD5, 0xDB)},
    {'text': '• 支持测试连接功能：输入域账号密码验证 LDAP 配置是否可用    • 配置优先级：数据库中的配置 > appsettings.json 文件配置', 'size': 13, 'color': RGBColor(0xD1, 0xD5, 0xDB)},
    {'text': '• 可配置项: 启用开关 / 服务器地址 / 端口 / SSL / 域名 / Base DN / 用户过滤器模板 / 服务账号凭据', 'size': 13, 'color': RGBColor(0xD1, 0xD5, 0xDB)},
], line_spacing=1.3)


# ============================================================
# SLIDE 11: DEPLOYMENT & HTTPS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, 'Docker 部署 & HTTPS 安全方案', 'DEPLOYMENT & HTTPS')
add_footer(slide, '11 / 14')

# Docker
add_rect(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(2.5), LIGHT)
add_text_box(slide, Inches(0.8), Inches(1.6), Inches(5.2), Inches(0.4),
             'Docker 容器化部署', font_size=18, color=PRIMARY, bold=True)
add_multi_text(slide, Inches(0.8), Inches(2.1), Inches(5.2), Inches(1.8), [
    'docker compose up -d  一键启动',
    '',
    '• API 容器 (excelimportsystem-api)',
    '  - .NET 10 ASP.NET Core 应用',
    '  - 端口映射: Host 5001 → Container 5000',
    '',
    '• Frontend 容器 (excelimportsystem-frontend)',
    '  - Nginx + React 静态文件',
    '  - HTTP :80 + HTTPS :443 双端口',
    '',
    '• SQL Server: 外部数据库',
    '  - 通过 host.docker.internal 连接宿主机',
], font_size=13, color=DARK, line_spacing=1.4)

# HTTPS
add_rect(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(2.5), LIGHT)
add_text_box(slide, Inches(7.1), Inches(1.6), Inches(5.2), Inches(0.4),
             'HTTPS 自建 CA 证书方案', font_size=18, color=SECONDARY, bold=True)
add_multi_text(slide, Inches(7.1), Inches(2.1), Inches(5.2), Inches(1.8), [
    '内部 CA 签发服务器证书，无需公网域名',
    '',
    '• ca.crt — CA 根证书 (有效期 10 年)',
    '• server.crt — 服务器证书 (CN=服务器IP)',
    '• HTTP → HTTPS 301 强制重定向',
    '• TLS 1.2 / 1.3, 禁用旧协议',
    '• HSTS Header (max-age=31536000)',
    '• 企业环境通过 GPO 自动分发 CA 证书',
], font_size=13, color=DARK, line_spacing=1.4)

# Offline deploy
add_text_box(slide, Inches(0.5), Inches(4.4), Inches(12), Inches(0.5),
             '离线部署包 (offline-deploy/)', font_size=18, color=DARK, bold=True)

deploy_items = [
    ('excelimportsystem-\napi.tar.gz', '后端镜像\n~135 MB', PRIMARY),
    ('excelimportsystem-\nfrontend.tar.gz', '前端镜像\n~32 MB', SECONDARY),
    ('docker-compose.yml', '编排文件\n(image: 模式)', ACCENT),
    ('generate-cert\n.sh / .bat', '证书生成\n脚本', ORANGE),
    ('nginx-https.conf', 'HTTPS\n配置模板', RGBColor(0x8B, 0x5C, 0xF6)),
    ('deploy.sh / .bat', '一键部署\n脚本', RGBColor(0xEC, 0x48, 0x99)),
]

for i, (name, desc, color) in enumerate(deploy_items):
    x = Inches(0.5 + i * 2.1)
    add_rect(slide, x, Inches(5.0), Inches(1.9), Inches(0.55), color)
    add_text_box(slide, x, Inches(5.02), Inches(1.9), Inches(0.5),
                 name, font_size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(5.6), Inches(1.9), Inches(0.5),
                 desc, font_size=10, color=GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 12: UI SHOWCASE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, '系统界面展示', 'UI SHOWCASE')
add_footer(slide, '12 / 14')

pages = [
    ('登录页面', PRIMARY, '验证码 + 双重认证\n安全防护提示\n中英文切换'),
    ('仪表板', SECONDARY, '统计卡片 (导入次数/行数)\nAdmin 看全部 / 用户看自己\n实时数据汇总'),
    ('数据导入', ACCENT, '服务器→数据库→表 三级选择\n文件拖拽上传\n列映射可视化\n实时进度轮询'),
    ('用户管理', ORANGE, 'CRUD 操作\nAD 关联/解除\n数据库权限配置\n密码重置'),
    ('服务器管理', RGBColor(0x8B, 0x5C, 0xF6), '多 SQL Server 注册\n连接测试\n启用/禁用开关\n编辑与删除'),
    ('导入日志 & 审计', RED, '导入记录查询 (表名/状态/日期)\n登录审计 (用户名/IP/状态/原因)\n分页浏览，详情感知'),
]

for i, (name, color, desc) in enumerate(pages):
    col = i % 3
    row = i // 3
    x = Inches(0.5 + col * 4.2)
    y = Inches(1.4 + row * 2.9)
    add_rect(slide, x, y, Inches(3.9), Inches(2.6), LIGHT)
    add_rect(slide, x, y, Inches(3.9), Inches(0.55), color)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.05), Inches(3.5), Inches(0.45),
                 name, font_size=17, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.75), Inches(3.3), Inches(1.6),
                 desc, font_size=14, color=DARK)


# ============================================================
# SLIDE 13: USE CASES
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, '应用场景', 'USE CASES & SCENARIOS')
add_footer(slide, '13 / 14')

scenarios = [
    ('🏭\n制造业', 'ERP/SCADA 数据入库', [
        '• 益海嘉里等企业 ERP 报表自动导入',
        '• 多工作表 Excel 自动识别',
        '• 空 Sheet 占位符智能跳过',
        '• 支持 .xls 旧格式 (97-2003)',
    ]),
    ('🏦\n金融保险', '业务数据批量导入', [
        '• 理赔/核保数据批量入库',
        '• 事务模式保证原子性',
        '• 完整操作日志可追溯',
        '• 表级权限隔离敏感数据',
    ]),
    ('🏢\n集团企业', '多子公司数据汇聚', [
        '• 多 SQL Server 实例统一管理',
        '• AD 域账号集成单点登录',
        '• RBAC 角色权限分级管控',
        '• HTTPS + CA 证书安全传输',
    ]),
    ('🏗️\n中小企业', '轻量级数据中台', [
        '• Docker Compose 一键部署',
        '• 离线部署包，无需公网',
        '• CSV/TSV/TXT 多种格式支持',
        '• 零代码添加新导入目标表',
    ]),
]

for i, (icon, title, items) in enumerate(scenarios):
    x = Inches(0.5 + i * 3.2)
    add_rect(slide, x, Inches(1.5), Inches(3.0), Inches(5.3), LIGHT)
    add_rect(slide, x, Inches(1.5), Inches(3.0), Inches(0.08), PRIMARY)
    add_text_box(slide, x + Inches(0.2), Inches(1.7), Inches(2.6), Inches(1.0),
                 icon, font_size=24, color=PRIMARY, bold=True)
    add_text_box(slide, x + Inches(0.2), Inches(2.7), Inches(2.6), Inches(0.4),
                 title, font_size=18, color=PRIMARY, bold=True)
    add_rect(slide, x + Inches(0.2), Inches(3.15), Inches(1.0), Inches(0.04), ACCENT)
    add_multi_text(slide, x + Inches(0.2), Inches(3.35), Inches(2.6), Inches(3.0),
                   items, font_size=13, color=DARK, line_spacing=1.8)


# ============================================================
# SLIDE 14: SUMMARY
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_rect(slide, Inches(0), Inches(0), SLD_W, Inches(0.08), PRIMARY)

add_text_box(slide, Inches(1.0), Inches(1.0), Inches(11), Inches(0.8),
             '总结与优势', font_size=40, color=WHITE, bold=True)

# Key advantages
advantages = [
    ('🚀', '开箱即用', 'Docker 一键部署，零代码添加新导入目标，30分钟完成从部署到首次导入'),
    ('🛡️', '安全合规', '10+ 层安全防护，完整审计日志，满足企业安全审计和合规要求'),
    ('⚡', '高性能', '异步 Fire-and-Forget + 批次处理，200MB 大文件稳定导入，实时进度反馈'),
    ('🏢', '企业级', 'LDAP/AD 集成，RBAC 权限控制，多服务器支持，HTTPS 内部 CA 证书'),
    ('🎨', '易用性', 'Web 化操作界面，中英文双语，智能列映射，拖拽上传，实时进度可视化'),
    ('🔧', '可扩展', 'Clean Architecture 架构清晰，动态表发现，新权限/新角色易于添加'),
]

for i, (icon, title, desc) in enumerate(advantages):
    col = i % 3
    row = i // 3
    x = Inches(0.5 + col * 4.2)
    y = Inches(2.1 + row * 2.4)
    add_rect(slide, x, y, Inches(3.9), Inches(2.1), RGBColor(0x15, 0x1F, 0x35))
    add_text_box(slide, x + Inches(0.2), y + Inches(0.15), Inches(3.5), Inches(0.5),
                 f'{icon}  {title}', font_size=20, color=ACCENT, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.75), Inches(3.3), Inches(1.2),
                 desc, font_size=13, color=RGBColor(0xD1, 0xD5, 0xDB))

# Contact info
add_rect(slide, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4), RGBColor(0x15, 0x1F, 0x35))
add_text_box(slide, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4),
             '技术栈: .NET 10 + React 19 + SQL Server  |  Clean Architecture  |  JWT + LDAP  |  Docker 容器化  |  Version 1.4.0',
             font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================
# SAVE
# ============================================================
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'Excel导入系统-推广介绍.pptx')
prs.save(output_path)
print(f'PPT saved to: {output_path}')
